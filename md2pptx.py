from markdown import Markdown
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Pt

class MarkdownHTMLParser(HTMLParser):
    """
    A custom HTML Parser for converting HTML derived from Markdown into a PowerPoint presentation using a provided template.
    Handles formatting for lists, bold text, and hyperlinks.
    """
    def __init__(self, presentation):
        super().__init__()
        self.presentation = presentation
        self.current_slide = None
        self.current_slide_index = 0
        self.title_placeholder = None
        self.text_placeholder = None
        self.current_paragraph = None
        self.in_list = False
        self.list_type = None
        self.list_counter = 0
        self.last_tag = None
        self.bold = False
        self.current_href = None  # To store URL for hyperlinks

    def add_or_duplicate_slide(self):
        if self.current_slide_index < len(self.presentation.slides):
            self.current_slide = self.presentation.slides[self.current_slide_index]
        else:
            last_layout = self.presentation.slides[-1].slide_layout
            self.current_slide = self.presentation.slides.add_slide(last_layout)
        self.title_placeholder = self.current_slide.shapes.title
        self.text_placeholder = self.current_slide.placeholders[1]
        self.current_slide_index += 1
        self.current_paragraph = None

    def handle_starttag(self, tag, attrs):
        self.last_tag = tag
        if tag in ["h1", "h2", "h3", "h4", "h5", "h6"]:
            if not self.current_slide or self.title_placeholder.text:
                self.add_or_duplicate_slide()
            self.current_paragraph = self.title_placeholder.text_frame.add_paragraph()
        elif tag in ["p", "li"]:
            if not self.current_slide:
                self.add_or_duplicate_slide()
            if not self.current_paragraph or tag == "li":
                self.current_paragraph = self.text_placeholder.text_frame.add_paragraph()
            if tag == "li":
                if self.list_type == "ol":
                    self.list_counter += 1
                    self.current_paragraph.text = f"{self.list_counter}. "
                elif self.list_type == "ul":
                    self.current_paragraph.text = "• "
        elif tag in ["ul", "ol"]:
            self.in_list = True
            self.list_type = tag
            self.list_counter = 0
        elif tag in ["strong", "b"]:
            self.bold = True
        elif tag == "a":
            self.current_href = dict(attrs).get("href", "")  # Extract href attribute

    def handle_endtag(self, tag):
        if tag in ["ul", "ol"]:
            self.in_list = False
            self.list_type = None
        elif tag in ["strong", "b"]:
            self.bold = False
        elif tag == "a":
            self.current_href = None  # Reset hyperlink after closing tag

    def handle_data(self, data):
        if self.current_paragraph:
            run = self.current_paragraph.add_run()
            run.text = data.strip()
            if self.bold:
                run.font.bold = True
            if self.current_href:  # Apply hyperlink if href is set
                run.hyperlink.address = self.current_href

def convert_markdown_to_pptx(markdown_text, template):
    markdown = Markdown(output_format='html')
    html_content = markdown.convert(markdown_text)
    prs = Presentation(template)
    parser = MarkdownHTMLParser(prs)
    parser.feed(html_content)
    prs.save('test.pptx')
    return 'test.pptx'